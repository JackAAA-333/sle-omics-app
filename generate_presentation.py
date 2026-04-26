from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

OUT = 'outputs/report_presentation.pptx'
REPORT = 'outputs/report.md'
FIG_DIRS = ['outputs/report_figs','outputs_advanced']

def add_title_slide(prs, title, subtitle=''):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_text_slide(prs, title, text):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    for line in text.split('\n')[:8]:
        p = body.add_paragraph() if body.text else body.paragraphs[0]
        p.text = line
        p.font.size = Pt(14)

def add_image_slide(prs, img_path, caption=''):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5); top = Inches(1)
    slide.shapes.add_picture(img_path, left, top, width=Inches(9))
    if caption:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.8), Inches(9), Inches(0.6))
        tf = tx.text_frame
        tf.text = caption

def main():
    prs = Presentation()
    add_title_slide(prs, 'SLE 多组学整合报告', '自动生成')
    if os.path.exists(REPORT):
        with open(REPORT,'r',encoding='utf-8') as f:
            txt = f.read()
        add_text_slide(prs, '方法与摘要', '\n'.join(txt.split('\n')[:12]))
    # add up to 6 figure slides
    count = 0
    for d in FIG_DIRS:
        if not os.path.isdir(d):
            continue
        for fn in sorted(os.listdir(d)):
            if count>=6:
                break
            if fn.lower().endswith(('.png','.jpg','.jpeg')):
                add_image_slide(prs, os.path.join(d,fn), caption=fn)
                count += 1
    prs.save(OUT)
    print('Presentation saved:', OUT)

if __name__=='__main__':
    main()
